from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # --- HELPER: Title Slide ---
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        title_placeholder.text = title
        subtitle_placeholder.text = subtitle

    # --- HELPER: Content Slide ---
    def add_content_slide(title, content_items, image_path=None):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        # If image, resize text to left
        if image_path:
             body_shape.width = Inches(5.5)
             body_shape.left = Inches(0.5)
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20) # Slightly smaller font to fit
            
        if image_path and os.path.exists(image_path):
            # Add image on the right
            slide.shapes.add_picture(image_path, Inches(6.2), Inches(2), width=Inches(3.5))


    # --- SLIDE 1: Link & Overview ---
    add_title_slide("ISAAC Project Update", "Schema, Portal, and Roadmap")

    # --- SLIDE 2: The Foundation (New) ---
    add_content_slide("1. The Foundation: AI-Ready Schema", [
        "• Based on our GitHub Wiki: 'The Semantic Layer above Raw Data'.",
        "• ISAAC Records capture meaning, not just bytes.",
        "• Key Components defined in the Wiki:",
        "   - Context: Physical system & conditions.",
        "   - Provenance: Origin of the information.",
        "   - Interpretation: Reduced metrics (e.g., efficiency).",
        "   - Pointers: Links to raw data files."
    ], image_path="screenshots/wiki_home.png")

    # --- SLIDE 3: The Portal ---
    add_content_slide("2. The New Ontology Portal", [
        "• We now have a live, dynamic website for the Schema.",
        "• Acts as the 'Single Source of Truth' for all terms.",
        "• 'What you see is what you document.'",
        "• Features a visual explorer (Concept Map) and navigation.",
        "• URL: http://localhost:8501 (Internal access only for now)"
    ], image_path="screenshots/ontology_map.png")
    
    # --- SLIDE 4: Schema Extensibility ---
    add_content_slide("3. The Schema Strategy: 'Big Blocks'", [
        "• The Schema is structured into core 'Super-Blocks':",
        "   - System, Sample, Context, Measurement, Descriptor",
        "• Approach: Add specificity *within* these blocks.",
        "• Example: Adding 'viscosity' to 'Context' -> 'context.transport.viscosity'.",
        "• Flexible: We can add properties without breaking the core structure."
    ])

    # --- SLIDE 5: The Validator ---
    add_content_slide("4. The Gatekeeper: Data Validator", [
        "• New quality control step BEFORE database ingestion.",
        "• Users upload their metadata (Excel) to the Portal.",
        "• Immediate feedback on:",
        "   - Required fields",
        "   - Correct terminology (from Ontology)",
        "   - Valid value ranges",
        "• Ensures only 'AI-Ready' records enter the system."
    ], image_path="screenshots/validator_page.png")

    # --- SLIDE 6: Roadmap ---
    add_content_slide("5. Future: Authentication & Identity", [
        "• Developing 'KID' (Kernel Identity) for robust identification.",
        "• Website Authentication:",
        "   - Moving from open access to role-based access.",
        "   - Controlling who can EDIT vs. VIEW the ontology.",
        "   - Secure API tokens for automated Agents.",
        "• Goal: A secure, trusted environment for scientific knowledge."
    ])

    # Save
    output_path = "ISAAC_All_Hands.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
